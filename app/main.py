from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pathlib import Path
import io
from pptx import Presentation
import re
from .models import RenderRequest
from .ppt import generate_presentation
import tempfile
import subprocess
import os

TEMPLATE_PATH = str(Path(__file__).parent / "templates" / "silicon_eic_template.pptx")

app = FastAPI()

# Ajusta orígenes si quieres restringirlo
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def _build_replacements(req: RenderRequest) -> dict:
    """
    Solo usamos los placeholders que tu plantilla necesita.
    Enviamos las llaves con {{...}} para que el replace sea exacto.
    """
    repl = {
        "{{COMPANY_NAME}}": req.company_name,
    }

    # pricing_overrides: solo añadimos los que vengan
    po = req.pricing_overrides or {}
    if po.SETUP_FEE is not None:
        repl["{{SETUP_FEE}}"] = str(po.SETUP_FEE)
    if po.SHORT_FEE is not None:
        repl["{{SHORT_FEE}}"] = str(po.SHORT_FEE)
    if po.FULL_FEE is not None:
        repl["{{FULL_FEE}}"] = str(po.FULL_FEE)
    if po.GRANT_FEE is not None:
        repl["{{GRANT_FEE}}"] = po.GRANT_FEE
    if po.EQUITY_FEE is not None:
        repl["{{EQUITY_FEE}}"] = po.EQUITY_FEE

    return repl


@app.post("/render")
async def render_presentation(request: RenderRequest):
    try:
        replacements = _build_replacements(request)
        buf = generate_presentation(
            template_path=TEMPLATE_PATH,
            replacements=replacements,
            slide_toggles=request.slide_toggles or {},
        )
        filename = f"proposal_{request.company_name.replace(' ', '_')}.pptx"
        return StreamingResponse(
            io.BytesIO(buf.getvalue()),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/render_pdf")
async def render_presentation_pdf(request: RenderRequest):
    """
    Nuevo endpoint que devuelve pdf en vez de pptx
    """
    try:
        if not os.path.exists(TEMPLATE_PATH):
            raise HTTPException(status_code=500, detail="Template file not found.")
        replacements = _build_replacements(request)
        if request.proposal_date:
            replacements["{{DATE}}"] = request.proposal_date.strftime("%B %d, %Y")
        else:
            replacements["{{DATE}}"] = ""

        pptx_stream = generate_presentation(
            template_path=TEMPLATE_PATH,
            replacements=replacements,
            slide_toggles=request.slide_toggles or {},
        )

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as pptx_temp:
            pptx_temp.write(pptx_stream.getvalue())
            pptx_temp.flush()
            pptx_path = pptx_temp.name

        pdf_path = pptx_path.replace(".pptx", ".pdf")

        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                pptx_path,
                "--outdir",
                os.path.dirname(pptx_path),
            ],
            check=True,
        )

        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

        os.remove(pptx_path)

        os.remove(pdf_path)

        safe_company = "".join(
            c if c.isalnum() or c in " _-" else "_" for c in request.company_name
        )
        filename = f"proposal_eic_template_{safe_company}.pdf"

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except subprocess.CalledProcessError as e:
        raise HTTPException(
            status_code=500, detail=f"error converting to libreoffice: {str(e)}"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"error general: {str(e)}")

@app.post("/render_pdf_custom")
async def render_pdf_custom(request: RenderRequest, remove_slides: list[int] = []):
    """
    Genera pdf desde la plantilla pptx, permitiendo eliminar diapositivas específicas antes de la conversión.
    """
    try:
        if not os.path.exists(TEMPLATE_PATH):
            raise HTTPException(status_code=500, detail="Template file not found.")

        # Construir los reemplazos
        replacements = _build_replacements(request)

        # Agregar fecha si está disponible
        if request.proposal_date:
            replacements["{{DATE}}"] = request.proposal_date.strftime("%B %d, %Y")
        else:
            replacements["{{DATE}}"] = ""

        # Convertir slides enviados por el usuario
        remove_idx = sorted([s - 1 for s in remove_slides if s > 0])

        # Generar la presentación pptx temporal
        pptx_stream = generate_presentation(
            template_path=TEMPLATE_PATH,
            replacements=replacements,
            slide_toggles=request.slide_toggles or {},
        )

        # Guardar temporalmente para poder manipularlo
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
            temp_pptx.write(pptx_stream.getvalue())
            temp_pptx.flush()
            pptx_path = temp_pptx.name

        # Abrir la presentación para eliminar las diapositivas especificadas
        from pptx import Presentation
        prs = Presentation(pptx_path)

        # Eliminar desde atrás hacia adelante
        for idx in reversed(remove_idx):
            if 0 <= idx < len(prs.slides._sldIdLst):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]

        # Guardar el pptx ya modificado
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as final_pptx:
            prs.save(final_pptx.name)
            final_pptx_path = final_pptx.name

        # Convertir a PDF usando libreoffice
        pdf_path = final_pptx_path.replace(".pptx", ".pdf")

        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                final_pptx_path,
                "--outdir", os.path.dirname(final_pptx_path),
            ],
            check=True,
        )

        # Leer el PDF generado
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

        # Limpiar archivos temporales
        os.remove(pptx_path)
        os.remove(final_pptx_path)
        os.remove(pdf_path)

        # Nombre final del archivo
        safe_company = "".join(
            c if c.isalnum() or c in " -" else "" for c in request.company_name
        )
        filename = f"proposal_custom_{safe_company}.pdf"

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except subprocess.CalledProcessError as e:
        raise HTTPException(
            status_code=500, detail=f"libreoffice conversion error: {str(e)}"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"General error: {str(e)}")
    
@app.post("/render_pptx_custom")
async def render_pptx_custom(request: RenderRequest, remove_slides: list[int] = []):
    """
    Igual que render_pdf_custom pero devolviendo PPTX modificado.
    """
    try:
        if not os.path.exists(TEMPLATE_PATH):
            raise HTTPException(status_code=500, detail="Template file not found.")

        # Build replacements
        replacements = _build_replacements(request)
        if request.proposal_date:
            replacements["{{DATE}}"] = request.proposal_date.strftime("%B %d, %Y")
        else:
            replacements["{{DATE}}"] = ""

        # Convert slide numbers (1-based → 0-based)
        remove_idx = sorted([s - 1 for s in remove_slides if s > 0])

        # 1. Create base PPTX
        pptx_stream = generate_presentation(
                template_path=TEMPLATE_PATH,
                replacements=replacements,
                slide_toggles=request.slide_toggles or {},
            )

        # 2. Save temp PPTX
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
            temp_pptx.write(pptx_stream.getvalue())
            temp_pptx.flush()
            pptx_path = temp_pptx.name

        # 3. Load PPTX and remove slides
        from pptx import Presentation
        prs = Presentation(pptx_path)

        for idx in reversed(remove_idx):
            if 0 <= idx < len(prs.slides._sldIdLst):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]

        # 4. Save final PPTX
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as final_pptx:
            prs.save(final_pptx.name)
            final_pptx_path = final_pptx.name

        # Read output file
        with open(final_pptx_path, "rb") as f:
            pptx_bytes = f.read()

        # Cleanup
        os.remove(pptx_path)
        os.remove(final_pptx_path)

        safe_company = "".join(
            c if c.isalnum() or c in " -" else "" for c in request.company_name
        )
        filename = f"proposal_custom_{safe_company}.pptx"

        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"General error: {str(e)}")